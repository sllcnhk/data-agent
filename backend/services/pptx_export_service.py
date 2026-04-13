"""
PPTX 导出服务

使用 Playwright 截图 + python-pptx 组装 PowerPoint 文件。

每个图表截图一页，首页封面带标题 + LLM 总结，末页附数据说明。

依赖：
  pip install playwright python-pptx Pillow
  playwright install chromium
"""
from __future__ import annotations

import asyncio
import io
import logging
import os
import tempfile
from pathlib import Path
from typing import List, Optional

logger = logging.getLogger(__name__)

# PPT 页面尺寸（宽屏 16:9）
_SLIDE_WIDTH_EMU = 12192000   # 33.87 cm
_SLIDE_HEIGHT_EMU = 6858000   # 19.05 cm


async def html_to_pptx(
    html_path: str,
    output_path: str,
    title: str = "数据分析报告",
    summary: str = "",
    wait_ms: int = 2000,
) -> str:
    """
    将 HTML 报告导出为 PPTX。

    流程：
      1. Playwright 打开 HTML，对每个 .chart-card 截图
      2. python-pptx 组装：封面 → 图表页 × N → 结语页

    Args:
        html_path:   HTML 文件绝对路径
        output_path: PPTX 输出路径
        title:       报告标题
        summary:     LLM 总结文字（放首页）
        wait_ms:     等待 ECharts 动画的毫秒数

    Returns:
        output_path
    """
    try:
        screenshots = await _take_screenshots(html_path, wait_ms)
    except ImportError:
        logger.error("[PPTX] playwright 未安装，无法生成 PPTX")
        raise RuntimeError("请先安装 playwright: pip install playwright && playwright install chromium")

    try:
        _build_pptx(output_path, title, summary, screenshots)
    except ImportError:
        logger.error("[PPTX] python-pptx 未安装，无法生成 PPTX")
        raise RuntimeError("请先安装 python-pptx: pip install python-pptx")

    logger.info("[PPTX] 导出完成: %s (%d 页)", output_path, len(screenshots) + 2)
    return output_path


async def _take_screenshots(html_path: str, wait_ms: int) -> List[bytes]:
    """使用 Playwright 对每个 chart-card 截图，返回 PNG bytes 列表。"""
    from playwright.async_api import async_playwright

    file_url = Path(html_path).as_uri()
    screenshots: List[bytes] = []

    async with async_playwright() as pw:
        browser = await pw.chromium.launch(
            headless=True,
            args=["--no-sandbox", "--disable-setuid-sandbox"],
        )
        page = await browser.new_page(viewport={"width": 1600, "height": 900})
        await page.goto(file_url, wait_until="networkidle", timeout=30000)
        await page.wait_for_timeout(wait_ms)

        # 先截全页图（作为概览幻灯片）
        full_shot = await page.screenshot(full_page=False, type="png")
        screenshots.append(full_shot)

        # 再逐个截 chart-card
        cards = await page.query_selector_all(".chart-card")
        for card in cards:
            try:
                shot = await card.screenshot(type="png")
                screenshots.append(shot)
            except Exception as e:
                logger.warning("[PPTX] 图表截图失败: %s", e)

        await browser.close()

    return screenshots


def _build_pptx(
    output_path: str,
    title: str,
    summary: str,
    screenshots: List[bytes],
) -> None:
    """用 python-pptx 组装幻灯片。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Emu(_SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(_SLIDE_HEIGHT_EMU)

    blank_layout = prs.slide_layouts[6]  # 全空白布局

    # ── 封面 ──────────────────────────────────────────────────────────────────
    cover = prs.slides.add_slide(blank_layout)
    _add_cover(cover, title, summary)

    # ── 全页概览（第一张截图）────────────────────────────────────────────────
    if screenshots:
        overview = prs.slides.add_slide(blank_layout)
        _add_screenshot_slide(overview, screenshots[0], "报告总览")

    # ── 每张图表独占一页 ──────────────────────────────────────────────────────
    for i, shot in enumerate(screenshots[1:], 1):
        slide = prs.slides.add_slide(blank_layout)
        _add_screenshot_slide(slide, shot, f"图表 {i}")

    # ── 末页 ──────────────────────────────────────────────────────────────────
    end_slide = prs.slides.add_slide(blank_layout)
    _add_end_slide(end_slide)

    prs.save(output_path)


def _add_cover(slide, title: str, summary: str) -> None:
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    # 背景渐变蓝（用矩形模拟）
    bg = slide.shapes.add_shape(1, 0, 0, Emu(_SLIDE_WIDTH_EMU), Emu(_SLIDE_HEIGHT_EMU))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x16, 0x77, 0xFF)
    bg.line.fill.background()

    # 白色装饰条
    deco = slide.shapes.add_shape(1, 0, Emu(int(_SLIDE_HEIGHT_EMU * 0.72)),
                                  Emu(_SLIDE_WIDTH_EMU), Emu(int(_SLIDE_HEIGHT_EMU * 0.28)))
    deco.fill.solid()
    deco.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    deco.line.fill.background()

    # 标题文字
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(14), Inches(1.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 总结文字
    if summary:
        summary_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(3.5), Inches(14), Inches(2.5))
        stf = summary_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = summary[:300] + ("…" if len(summary) > 300 else "")
        sp.font.size = Pt(14)
        sp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 日期
    from datetime import datetime
    date_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(6), Inches(0.5))
    dp = date_box.text_frame.paragraphs[0]
    dp.text = datetime.now().strftime("%Y-%m-%d 生成")
    dp.font.size = Pt(12)
    dp.font.color.rgb = RGBColor(0x59, 0x59, 0x59)


def _add_screenshot_slide(slide, shot_bytes: bytes, label: str) -> None:
    from pptx.util import Inches, Pt, Emu
    import io

    # 将图片铺满幻灯片（留少量边距）
    img_stream = io.BytesIO(shot_bytes)
    pic = slide.shapes.add_picture(
        img_stream,
        left=Emu(int(_SLIDE_WIDTH_EMU * 0.01)),
        top=Emu(int(_SLIDE_HEIGHT_EMU * 0.06)),
        width=Emu(int(_SLIDE_WIDTH_EMU * 0.98)),
        height=Emu(int(_SLIDE_HEIGHT_EMU * 0.92)),
    )


def _add_end_slide(slide) -> None:
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    bg = slide.shapes.add_shape(1, 0, 0, Emu(_SLIDE_WIDTH_EMU), Emu(_SLIDE_HEIGHT_EMU))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xFA)
    bg.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(3), Inches(2.5), Inches(8), Inches(2))
    p = tx.text_frame.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x16, 0x77, 0xFF)

    sub = tx.text_frame.add_paragraph()
    sub.text = "数据驱动决策"
    sub.font.size = Pt(18)
    sub.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
